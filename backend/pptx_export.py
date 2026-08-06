"""Native PowerPoint export — real text runs, not a picture of a slide.

This is the capability the whole structured-master architecture exists to buy.
Claude Design emits HTML and every other format is a translation away from it,
so its PPTX export flattens text to images, substitutes fonts and drops master
slides; the recurring complaint is that the result is "a beautiful slide nobody
can edit". Create EQ's own PDF path has the same shape of problem for a
different reason — it rasterises the DOM through html2canvas.

Here the deck is already a structured element tree (x/y/w/h + typography +
colour), which maps almost directly onto DrawingML. Every text block becomes a
real PowerPoint text frame: selectable, restyleable, translatable, and editable
by a sales rep who has never opened this product.

Geometry note: element coordinates are CSS pixels at 96 DPI, and PowerPoint
measures in EMU at 914400 per inch, so the conversion is exactly 9525 EMU per
pixel. That is a lossless integer scale, which is why the exported slide lands
on the same grid as the canvas rather than drifting.
"""

import base64
import io
import logging
import re
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

log = logging.getLogger(__name__)

EMU_PER_PX = 9525  # 914400 EMU/inch ÷ 96 px/inch — exact, no rounding drift

ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT,
         "justify": PP_ALIGN.JUSTIFY}

# Fallback palette for decks that still store semantic colour tokens rather than
# literal hex. Premium/Design EQ decks emit hex directly (that is what lets one
# deck mix light and dark surfaces), so this only matters for older decks.
DEFAULT_PALETTE = {"bg": "#0F1010", "bg2": "#1F2937", "accent": "#E85D3A",
                   "text": "#FAFAFA", "muted": "#9CA3AF"}

_HEX_RE = re.compile(r"^#?([0-9a-fA-F]{3}|[0-9a-fA-F]{6})$")


def px(v: float) -> Emu:
    return Emu(int(round(float(v or 0) * EMU_PER_PX)))


def _hex_to_rgb(value: str) -> Optional[RGBColor]:
    if not isinstance(value, str):
        return None
    m = _HEX_RE.match(value.strip())
    if not m:
        return None
    h = m.group(1)
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def resolve_color(value: Any, palette: Dict[str, str]) -> Optional[RGBColor]:
    """Mirrors the frontend's resolveColor: a palette token resolves against the
    deck palette, anything else is treated as a literal colour."""
    if not isinstance(value, str) or not value:
        return None
    if value in palette:
        return _hex_to_rgb(palette[value])
    return _hex_to_rgb(value)


def _palette_for(doc: Dict[str, Any]) -> Dict[str, str]:
    ai = doc.get("ai_palette")
    if isinstance(ai, dict) and ai.get("bg"):
        return {k: v for k, v in ai.items() if isinstance(v, str)}
    return dict(DEFAULT_PALETTE)


def _set_char_props(run, *, tracking_em: Optional[float], size_pt: float, uppercase: bool) -> None:
    """Letter-spacing and all-caps have no python-pptx API, so they are written
    straight onto the run's rPr. Both matter here: the design system's eyebrows
    are tracked mono small-caps, and dropping either turns them into ordinary
    body text. `spc` is in hundredths of a point; `cap` preserves the original
    string so the text stays editable rather than being upper-cased in place."""
    rPr = run._r.get_or_add_rPr()
    if tracking_em:
        rPr.set("spc", str(int(round(tracking_em * size_pt * 100))))
    if uppercase:
        rPr.set("cap", "all")


def _add_text(slide, el: Dict[str, Any], palette: Dict[str, str]) -> None:
    box = slide.shapes.add_textbox(px(el.get("x")), px(el.get("y")),
                                   px(el.get("w")), px(el.get("h")))
    tf = box.text_frame
    tf.word_wrap = True
    # The canvas measures a text block from its top-left with no padding; the
    # PowerPoint default insets would shift every block down and right.
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    size_pt = float(el.get("size") or 24) * 0.75  # CSS px -> points
    lines = str(el.get("text") or "").split("\n")
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = ALIGN.get(el.get("align") or "left", PP_ALIGN.LEFT)
        lh = el.get("line_height")
        if isinstance(lh, (int, float)) and lh > 0:
            para.line_spacing = float(lh)
        run = para.add_run()
        run.text = line
        f = run.font
        f.name = el.get("font") or "Geist"
        f.size = Pt(size_pt)
        weight = el.get("weight")
        f.bold = bool(isinstance(weight, (int, float)) and weight >= 600) or bool(el.get("bold"))
        f.italic = bool(el.get("italic"))
        color = resolve_color(el.get("color"), palette)
        if color is not None:
            f.color.rgb = color
        _set_char_props(run, tracking_em=el.get("letter_spacing"),
                        size_pt=size_pt, uppercase=bool(el.get("uppercase")))


def _add_shape(slide, el: Dict[str, Any], palette: Dict[str, str]) -> None:
    kind = el.get("shape")
    w, h = float(el.get("w") or 0), float(el.get("h") or 0)
    if w <= 0 or h <= 0:
        return
    radius = float(el.get("radius") or 0)
    if kind == "circle" or radius >= min(w, h) / 2:
        auto = MSO_SHAPE.OVAL
    elif radius > 0:
        auto = MSO_SHAPE.ROUNDED_RECTANGLE
    else:
        auto = MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(auto, px(el.get("x")), px(el.get("y")), px(w), px(h))
    fill_col = resolve_color(el.get("fill"), palette)
    if fill_col is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_col
    else:
        shp.fill.background()
    shp.line.fill.background()
    shp.shadow.inherit = False
    opacity = el.get("opacity")
    if isinstance(opacity, (int, float)) and 0 <= opacity < 1 and fill_col is not None:
        # Alpha is not exposed by python-pptx; written onto the solidFill so a
        # deliberately translucent panel does not export as a solid block.
        srgb = shp.fill.fore_color._xFill.find(qn("a:srgbClr"))
        if srgb is not None:
            alpha = srgb.makeelement(qn("a:alpha"), {"val": str(int(opacity * 100000))})
            srgb.append(alpha)


def _add_image(slide, el: Dict[str, Any]) -> None:
    src = el.get("src")
    if not isinstance(src, str) or not src:
        return
    try:
        if src.startswith("data:"):
            b64 = src.split(",", 1)[1]
            stream = io.BytesIO(base64.b64decode(b64))
        else:
            # Remote fetch is deliberately skipped: a slow or dead asset host
            # would stall the whole export, and a missing picture is a better
            # outcome than a request that never returns.
            return
        slide.shapes.add_picture(stream, px(el.get("x")), px(el.get("y")),
                                 px(el.get("w")), px(el.get("h")))
    except Exception as ex:
        log.warning("pptx export skipped an image: %s", ex)


def _paint_background(slide, bg: Any, palette: Dict[str, str], canvas: Dict[str, int]) -> None:
    """A full-bleed rectangle rather than the slide's own background fill —
    gradients and image backgrounds degrade to their base colour this way, and
    it keeps a single code path for every background type."""
    color = None
    if isinstance(bg, dict):
        color = resolve_color(bg.get("color") or bg.get("from") or "bg", palette)
    if color is None:
        color = resolve_color("bg", palette)
    if color is None:
        return
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0),
                                 px(canvas["w"]), px(canvas["h"]))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False


def deck_to_pptx(doc: Dict[str, Any]) -> bytes:
    """Convert a stored deck (Create EQ carousel shape) into a .pptx file.

    Returns the raw bytes. Every text element lands as an editable text frame;
    nothing is rasterised.
    """
    canvas = doc.get("canvas") or {"w": 1080, "h": 1350}
    palette = _palette_for(doc)

    prs = Presentation()
    prs.slide_width = px(canvas["w"])
    prs.slide_height = px(canvas["h"])
    blank = prs.slide_layouts[6]  # the only layout with no placeholders

    for s in doc.get("slides") or []:
        slide = prs.slides.add_slide(blank)
        _paint_background(slide, s.get("bg"), palette, canvas)
        # Elements are painted in array order, which is the deck's z-order.
        for el in s.get("elements") or []:
            try:
                t = el.get("type")
                if t == "text":
                    _add_text(slide, el, palette)
                elif t in ("shape", "line"):
                    _add_shape(slide, el, palette)
                elif t == "image":
                    _add_image(slide, el)
                # Charts, badges and icons have no faithful DrawingML analogue
                # yet; skipping one element is better than corrupting the file.
            except Exception as ex:
                log.warning("pptx export skipped an element (%s): %s", el.get("type"), ex)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def deck_text_stats(doc: Dict[str, Any]) -> Dict[str, int]:
    """Counts used to prove the export is live text rather than pictures."""
    slides = doc.get("slides") or []
    els: List[Dict[str, Any]] = [e for s in slides for e in (s.get("elements") or [])]
    return {
        "slides": len(slides),
        "text_elements": sum(1 for e in els if e.get("type") == "text"),
        "shape_elements": sum(1 for e in els if e.get("type") in ("shape", "line")),
        "image_elements": sum(1 for e in els if e.get("type") == "image"),
    }
