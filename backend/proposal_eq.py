"""Proposal EQ — AI research + proposal generation agent.

Fifth agent in the Innoira Agentic Suite: researches a lead (CRM activity
timeline + live web search) and drafts a branded proposal deck in the same
slide/element shape Create EQ uses, exportable to PDF (client-side, reusing
Create EQ's render pipeline) and PPTX (server-side, python-pptx).
"""

import io
import json
import uuid
from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from typing import List, Optional, Dict, Any

import anthropic

from server import (
    db, current_user, now_iso, new_id, _audit, _log_activity,
    _llm_chat, _extract_json, ANTHROPIC_API_KEY, ANTHROPIC_MODEL, STAGES,
)

proposal_router = APIRouter(prefix="/proposal-eq")


# ----------------------------- Models ------------------------------------------
class PricingItemIn(BaseModel):
    name: str
    price: str
    unit: str = ""
    description: str = ""


class ProposalGenIn(BaseModel):
    lead_id: str
    topic: str = ""
    include_pricing: bool = True


# ----------------------------- Pricing catalog ----------------------------------
@proposal_router.get("/pricing-catalog")
async def list_pricing(user=Depends(current_user)):
    return await db.pricing_catalog.find({"workspace_id": user["workspace_id"]}, {"_id": 0}).to_list(200)


@proposal_router.post("/pricing-catalog")
async def create_pricing_item(body: PricingItemIn, user=Depends(current_user)):
    doc = body.model_dump()
    doc.update({"id": new_id(), "workspace_id": user["workspace_id"], "created_at": now_iso()})
    await db.pricing_catalog.insert_one(doc)
    doc.pop("_id", None)
    return doc


@proposal_router.delete("/pricing-catalog/{item_id}")
async def delete_pricing_item(item_id: str, user=Depends(current_user)):
    await db.pricing_catalog.delete_one({"id": item_id, "workspace_id": user["workspace_id"]})
    return {"ok": True}


# ----------------------------- Research + draft ------------------------------------
async def _research_and_draft(lead: Dict[str, Any], deal: Optional[Dict[str, Any]], topic: str,
                               timeline: List[Dict[str, Any]]) -> tuple:
    """Single Claude call that both researches the company (web_search tool, only
    when a company name exists) and drafts the proposal JSON — merged into one
    request instead of two back-to-back calls, since two heavy sequential calls
    reliably tripped the same per-minute token rate limit in testing."""
    company = lead.get("company", "")
    if not ANTHROPIC_API_KEY:
        return "", {}
    system = (
        "You are Proposal EQ. "
        + ("If useful, briefly web-search the lead's company for 1-2 current, specific facts "
           "(recent news, scale, focus) before answering. " if company else "")
        + "Then respond with ONLY a STRICT JSON object (no other text, no markdown fences), schema: "
        '{"research_summary": str, "proposal": {'
        '"cover":{"title":str,"subtitle":str},'
        '"problem":{"title":str,"body":str},'
        '"solution":{"title":str,"body":str},'
        '"case_study":{"title":str,"body":str},'
        '"next_steps":{"title":str,"body":str,"cta":str}}} '
        "research_summary: 3-4 plain sentences on the company + a plausible pain point, or \"\" if no company given. "
        "proposal titles <=8 words, bodies <=55 words, plain text, no emojis. "
        "Use the CRM history to personalize the problem/solution — reference specifics when relevant. "
        "Omit case_study (empty strings) if there's nothing credible to say."
    )
    user_text = json.dumps({
        "lead": {"name": f"{lead.get('first_name', '')} {lead.get('last_name', '')}".strip(),
                 "company": company, "title": lead.get("title")},
        "deal": {"title": deal.get("title"), "stage": deal.get("stage")} if deal else None,
        "crm_history": timeline,
        "context": topic,
    })
    try:
        kwargs: Dict[str, Any] = dict(
            model=ANTHROPIC_MODEL, max_tokens=1536, system=system,
            messages=[{"role": "user", "content": user_text}],
        )
        if company:
            kwargs["tools"] = [{"type": "web_search_20250305", "name": "web_search", "max_uses": 2}]
        client = anthropic.AsyncAnthropic(api_key=ANTHROPIC_API_KEY)
        resp = await client.messages.create(**kwargs)
        text = "".join(b.text for b in resp.content if getattr(b, "type", None) == "text")
        parsed = _extract_json(text) or {}
        return parsed.get("research_summary", ""), parsed.get("proposal", {})
    except Exception:
        return "", {}


# ----------------------------- Slide building (mirrors legacySlideToElements) ----
def _elid() -> str:
    return str(uuid.uuid4())


def _text_el(x, y, w, h, text, font, size, weight, color, **extra) -> Dict[str, Any]:
    return {"id": _elid(), "type": "text", "x": x, "y": y, "w": w, "h": h, "text": text,
            "font": font, "size": size, "weight": weight, "color": color, **extra}


def _badge_el(x, y, text, **extra) -> Dict[str, Any]:
    return {"id": _elid(), "type": "badge", "x": x, "y": y, "text": text,
            "bg": "accent", "color": "bg", "radius": 999, "size": 22, **extra}


def _build_slide(subtitle: str, title: str, body: str, cta: Optional[str] = None) -> Dict[str, Any]:
    els: List[Dict[str, Any]] = []
    if subtitle:
        els.append(_text_el(80, 120, 920, 60, subtitle, "JetBrains Mono", 22, 500, "muted",
                             uppercase=True, letter_spacing=0.2, align="left"))
    if title:
        els.append(_text_el(80, 220, 920, 420, title, "Archivo Black", 84, 900, "accent",
                             line_height=0.98, align="left"))
    if body:
        els.append(_text_el(80, 680, 920, 560, body, "Inter", 28, 400, "text",
                             line_height=1.5, align="left"))
    if cta:
        els.append(_badge_el(80, 1220, cta))
    return {"_k": new_id(), "bg": {"type": "solid", "color": "bg"}, "elements": els}


def _fallback_content(lead: Dict[str, Any]) -> Dict[str, Any]:
    company = lead.get("company") or lead.get("first_name") or "your team"
    return {
        "cover": {"title": f"Proposal for {company}", "subtitle": "Prepared by Innoira"},
        "problem": {"title": "The Challenge", "body": "Add the prospect's core problem here."},
        "solution": {"title": "Our Approach", "body": "Describe the proposed solution here."},
        "case_study": {"title": "", "body": ""},
        "next_steps": {"title": "Let's talk", "body": "Ready to move forward.", "cta": "Book a call"},
    }


def _build_deck(lead: Dict[str, Any], content: Dict[str, Any],
                 pricing_items: List[Dict[str, Any]], include_pricing: bool) -> List[Dict[str, Any]]:
    if not content:
        content = _fallback_content(lead)

    slides = []
    cov = content.get("cover", {})
    slides.append(_build_slide(cov.get("subtitle", "Proposal"),
                                cov.get("title", f"For {lead.get('company') or lead.get('first_name')}"), ""))
    prob = content.get("problem", {})
    slides.append(_build_slide("The Challenge", prob.get("title", ""), prob.get("body", "")))
    sol = content.get("solution", {})
    slides.append(_build_slide("Our Approach", sol.get("title", ""), sol.get("body", "")))
    if include_pricing and pricing_items:
        lines = "\n".join(f"{p['name']} — {p['price']}{('/' + p['unit']) if p.get('unit') else ''}" for p in pricing_items)
        slides.append(_build_slide("Investment", "Pricing", lines))
    cs = content.get("case_study", {})
    if cs.get("body"):
        slides.append(_build_slide("Proof", cs.get("title") or "Case Study", cs.get("body", "")))
    nxt = content.get("next_steps", {})
    slides.append(_build_slide("Next Steps", nxt.get("title", "Let's talk"), nxt.get("body", ""),
                                cta=nxt.get("cta", "Book a call")))
    return slides


# ----------------------------- Proposals ------------------------------------------
@proposal_router.get("/proposals")
async def list_proposals(user=Depends(current_user)):
    items = await db.proposals.find({"workspace_id": user["workspace_id"]}, {"_id": 0}).sort("created_at", -1).to_list(200)
    for p in items:
        p["lead"] = await db.leads.find_one(
            {"id": p["lead_id"]}, {"_id": 0, "first_name": 1, "last_name": 1, "company": 1})
    return items


@proposal_router.get("/proposals/{pid}")
async def get_proposal(pid: str, user=Depends(current_user)):
    p = await db.proposals.find_one({"id": pid, "workspace_id": user["workspace_id"]}, {"_id": 0})
    if not p:
        raise HTTPException(404, "not found")
    return p


@proposal_router.post("/generate")
async def generate_proposal(body: ProposalGenIn, user=Depends(current_user)):
    lead = await db.leads.find_one({"id": body.lead_id, "workspace_id": user["workspace_id"]}, {"_id": 0})
    if not lead:
        raise HTTPException(404, "lead not found")
    deal = await db.deals.find_one({"lead_id": lead["id"], "workspace_id": user["workspace_id"]}, {"_id": 0})
    timeline_raw = await db.activities.find(
        {"lead_id": lead["id"], "workspace_id": user["workspace_id"]}, {"_id": 0}
    ).sort("at", -1).to_list(20)
    timeline = [{"type": a["type"], "summary": a["summary"]} for a in timeline_raw]
    pricing_items = await db.pricing_catalog.find({"workspace_id": user["workspace_id"]}, {"_id": 0}).to_list(50)

    research_text, content = await _research_and_draft(lead, deal, body.topic, timeline)
    slides = _build_deck(lead, content, pricing_items, body.include_pricing)

    doc = {
        "id": new_id(), "workspace_id": user["workspace_id"], "owner_id": user["id"],
        "lead_id": lead["id"], "deal_id": deal["id"] if deal else None,
        "topic": body.topic or f"Proposal for {lead.get('company') or lead.get('first_name')}",
        "status": "draft", "slides": slides, "research_notes": research_text,
        "palette_id": "midnight", "brand": {"logo_url": None, "colors": [], "fonts": []},
        "created_at": now_iso(), "updated_at": now_iso(), "sent_at": None,
    }
    await db.proposals.insert_one(doc)
    doc.pop("_id", None)
    await _audit(user, "proposal_eq.proposal.generate", {"id": doc["id"], "lead_id": lead["id"]})
    await _log_activity(user["workspace_id"], lead["id"], "proposal", "proposal_generated",
                         f"Generated proposal: {doc['topic']}", {"proposal_id": doc["id"]})
    return doc


@proposal_router.put("/proposals/{pid}")
async def update_proposal(pid: str, body: Dict[str, Any], user=Depends(current_user)):
    allowed = {k: v for k, v in body.items() if k in {"slides", "status", "topic", "palette_id", "brand"}}
    allowed["updated_at"] = now_iso()
    await db.proposals.update_one({"id": pid, "workspace_id": user["workspace_id"]}, {"$set": allowed})
    p = await db.proposals.find_one({"id": pid, "workspace_id": user["workspace_id"]}, {"_id": 0})
    if not p:
        raise HTTPException(404, "not found")
    return p


@proposal_router.post("/proposals/{pid}/mark-sent")
async def mark_proposal_sent(pid: str, user=Depends(current_user)):
    p = await db.proposals.find_one({"id": pid, "workspace_id": user["workspace_id"]}, {"_id": 0})
    if not p:
        raise HTTPException(404, "not found")
    await db.proposals.update_one({"id": pid}, {"$set": {"status": "sent", "sent_at": now_iso()}})
    await _log_activity(user["workspace_id"], p["lead_id"], "proposal", "proposal_sent",
                         f"Sent proposal “{p['topic']}”", {"proposal_id": pid})
    existing_deal = await db.deals.find_one({"lead_id": p["lead_id"], "workspace_id": user["workspace_id"]}, {"_id": 0})
    if existing_deal and existing_deal.get("stage") in ("new", "qualified", "meeting"):
        await db.deals.update_one({"id": existing_deal["id"]}, {"$set": {"stage": "proposal"}})
    elif not existing_deal:
        await db.deals.insert_one({
            "id": new_id(), "workspace_id": user["workspace_id"], "lead_id": p["lead_id"],
            "title": p["topic"], "value": 5000, "stage": "proposal", "created_at": now_iso(),
            "source_proposal_id": pid,
        })
    return {"ok": True}


@proposal_router.delete("/proposals/{pid}")
async def delete_proposal(pid: str, user=Depends(current_user)):
    await db.proposals.delete_one({"id": pid, "workspace_id": user["workspace_id"]})
    return {"ok": True}


# ----------------------------- PPTX export ------------------------------------------
_PALETTE_HEX = {
    "midnight": {"bg": "0F1010", "accent": "E85D3A", "text": "FAFAFA", "muted": "9CA3AF"},
    "bone":     {"bg": "E8E9EB", "accent": "212025", "text": "0F1010", "muted": "525252"},
    "sunset":   {"bg": "FF6B4A", "accent": "0F172A", "text": "FFFFFF", "muted": "FCD34D"},
    "ocean":    {"bg": "0A2540", "accent": "22D3EE", "text": "F0F9FF", "muted": "7DD3FC"},
    "forest":   {"bg": "14532D", "accent": "FCD34D", "text": "F0FDF4", "muted": "86EFAC"},
    "rose":     {"bg": "831843", "accent": "F9A8D4", "text": "FFF1F2", "muted": "FBCFE8"},
    "paper":    {"bg": "F5F1E8", "accent": "B45309", "text": "1C1917", "muted": "78716C"},
    "cyber":    {"bg": "030712", "accent": "34D399", "text": "F9FAFB", "muted": "4ADE80"},
    "coral":    {"bg": "FEE2E2", "accent": "DC2626", "text": "7F1D1D", "muted": "F97316"},
    "mono":     {"bg": "FFFFFF", "accent": "000000", "text": "000000", "muted": "71717A"},
}
_PX_TO_EMU = 9525  # 96 DPI: 914400 EMU per inch / 96 px per inch


@proposal_router.get("/proposals/{pid}/export.pptx")
async def export_pptx(pid: str, user=Depends(current_user)):
    p = await db.proposals.find_one({"id": pid, "workspace_id": user["workspace_id"]}, {"_id": 0})
    if not p:
        raise HTTPException(404, "not found")

    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor

    pal = _PALETTE_HEX.get(p.get("palette_id", "midnight"), _PALETTE_HEX["midnight"])

    def hexc(key: str) -> RGBColor:
        return RGBColor.from_string(pal.get(key, pal["text"]))

    prs = Presentation()
    prs.slide_width = Emu(1080 * _PX_TO_EMU)
    prs.slide_height = Emu(1350 * _PX_TO_EMU)
    blank_layout = prs.slide_layouts[6]

    for slide_data in p.get("slides", []):
        slide = prs.slides.add_slide(blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = hexc("bg")
        for el in slide_data.get("elements", []):
            if el.get("type") not in ("text", "badge"):
                continue
            x, y = el.get("x", 0), el.get("y", 0)
            w, h = el.get("w", 800), el.get("h", 100)
            box = slide.shapes.add_textbox(Emu(x * _PX_TO_EMU), Emu(y * _PX_TO_EMU),
                                            Emu(w * _PX_TO_EMU), Emu(h * _PX_TO_EMU))
            tf = box.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            run = para.add_run()
            run.text = el.get("text", "")
            run.font.size = Pt(max(8, int(el.get("size", 24) * 0.62)))
            run.font.bold = el.get("weight", 400) >= 700
            run.font.color.rgb = hexc(el.get("color", "text"))

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    filename = f"{(p.get('topic') or 'proposal')[:40].replace(' ', '-')}.pptx"
    return StreamingResponse(
        buf, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )
